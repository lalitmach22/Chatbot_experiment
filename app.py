import os
import streamlit as st
import re
import json
from datetime import datetime, timedelta
import pytz
from docx import Document  # For .docx
import pandas as pd
from bs4 import BeautifulSoup
import json
import yaml
from pptx import Presentation
from zipfile import ZipFile
import mimetypes 
import tempfile
from langchain_groq import ChatGroq
from langchain_community.document_loaders import PyPDFLoader
from langchain.chains import ConversationalRetrievalChain
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from supabase import create_client, Client

# Supabase credentials
url = os.getenv("url")
key = os.getenv("key")
supabase: Client = create_client(url, key)

GROQ_API_KEY = os.getenv("GROQ_API_KEY")

# Load the model
@st.cache_resource
def load_model():
    return ChatGroq(temperature=0.8, model="llama3-8b-8192")

def clean_text(text):
    # Replace multiple spaces with a single space
    text = re.sub(r'\s+', ' ', text)
    # Fix words broken by line breaks or formatting
    # Matches a lowercase/uppercase word followed by a line break without a space
    text = re.sub(r'(?<=[a-zA-Z])(?=[A-Z])', ' ', text)  # Add space between words when they are stuck together    
    # Fix punctuation followed by words with no space
    text = re.sub(r'(?<=[.,?!;])(?=[a-zA-Z])', ' ', text)  # Add space after punctuation if missing
    # Standardize newlines for better formatting
    text = re.sub(r'\.\s+', '.\n', text)  # Add newlines after sentences
    text = re.sub(r'(?<=:)\s+', '\n', text)  # Add newlines after colons
    # Additional cleanup (if needed)    text = text.strip()  # Remove leading and trailing whitespace
    return text

@st.cache_data
def load_hidden_documents(directory="hidden_docs"):
    """Load all supported file types from a directory and return their content."""
    all_texts = []

    for filename in os.listdir(directory):
        file_path = os.path.join(directory, filename)
        mime_type, _ = mimetypes.guess_type(file_path)

        try:
            # Handle PDF files
            if filename.endswith(".pdf"):
                loader = PyPDFLoader(file_path)
                pages = loader.load_and_split()
                all_texts.extend([page.page_content for page in pages])

            # Handle Word files (.docx)
            elif filename.endswith(".docx"):
                doc = Document(file_path)
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                all_texts.append(text)

            # Handle Text files (.txt)
            elif filename.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

            # Handle Excel files (.xlsx and .xls)
            elif filename.endswith(('.xlsx', '.xls')):
                excel_data = pd.read_excel(file_path)
                text = excel_data.to_string(index=False)
                all_texts.append(text)

            # Handle CSV files (.csv)
            elif filename.endswith(".csv"):
                csv_data = pd.read_csv(file_path)
                text = csv_data.to_string(index=False)
                all_texts.append(text)

            # Handle Markdown files (.md)
            elif filename.endswith(".md"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

            # Handle HTML files (.html, .htm)
            elif filename.endswith(('.html', '.htm')):
                with open(file_path, "r", encoding="utf-8") as file:
                    soup = BeautifulSoup(file, "html.parser")
                    all_texts.append(soup.get_text())

            # Handle JSON files (.json)
            elif filename.endswith(".json"):
                with open(file_path, "r", encoding="utf-8") as file:
                    data = json.load(file)
                    all_texts.append(json.dumps(data, indent=2))

            # Handle YAML files (.yaml, .yml)
            elif filename.endswith(('.yaml', '.yml')):
                with open(file_path, "r", encoding="utf-8") as file:
                    data = yaml.safe_load(file)
                    all_texts.append(json.dumps(data, indent=2))

            # Handle PowerPoint files (.pptx)
            elif filename.endswith(".pptx"):
                presentation = Presentation(file_path)
                for slide in presentation.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            slide_text.append(shape.text)
                    all_texts.append("\n".join(slide_text))

            # Handle ZIP files (.zip)
            elif filename.endswith(".zip"):
                with ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall("temp_extracted")
                    all_texts.extend(load_hidden_documents("temp_extracted"))

            # Handle Log files (.log)
            elif filename.endswith(".log"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

            # Handle unknown file types (fallback to text-based reading)
            elif mime_type and mime_type.startswith("text"):
                with open(file_path, "r", encoding="utf-8") as file:
                    all_texts.append(file.read())

        except Exception as e:
            print(f"Failed to process {filename}: {e}")
    cleaned_texts = [clean_text(text) for text in all_texts]
    return cleaned_texts

@st.cache_data
def save_to_supabase(all_texts):
    """Save the list of documents to the Supabase 'all_texts' table."""
    for text in all_texts:
        data = {"all_texts": text}
        response = supabase.table("all_texts").insert(data).execute()        
        # Check the response for success or failure
        if response.data:  # If the response contains data, the insert was successful
            print(f"Successfully saved: {text[:30]}...")
        else:  # If there is an error
            print(f"Failed to save text")
            
# Load documents and save to Supabase
all_texts = load_hidden_documents()
save_to_supabase(all_texts)

# Create vector store
@st.cache_resource
def create_vector_store(document_texts):
    embedder = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_texts(document_texts, embedder)

# Get modification times for documents
@st.cache_data
def get_file_mod_times(directory):
    """Get the modification times of all files in the directory."""
    return {
        f: os.path.getmtime(os.path.join(directory, f))
        for f in os.listdir(directory)
        if os.path.isfile(os.path.join(directory, f))  # Ensure it's a file, not a directory
    }

# Reload vector store if needed
def reload_vector_store_if_needed():
    current_mod_times = get_file_mod_times("hidden_docs")

    # Check if file modifications have occurred
    if "file_mod_times" not in st.session_state or st.session_state["file_mod_times"] != current_mod_times:
        st.session_state["file_mod_times"] = current_mod_times
        document_texts = load_hidden_documents()
        vector_store = create_vector_store(document_texts)  # Create vector store
        st.session_state["vector_store"] = vector_store    # Save in session state
    else:
        # Retrieve from session state if already initialized
        vector_store = st.session_state.get("vector_store", None)

    # Return the vector store (even if None)
    return vector_store

# Load model and vector store
model = load_model()

# Initialize vector_store
vector_store = reload_vector_store_if_needed()

# If still None, raise an error to debug initialization
if vector_store is None:
    raise ValueError("Failed to initialize vector_store. Ensure hidden_docs folder and embeddings setup are correct.")
##### my email lalitmach22@gmail.com
# Validate email
def is_valid_email(email):
    email_regex = r"^\d{2}f\d{7}@ds\.study\.iitm\.ac\.in$"
    return re.match(email_regex, email) is not None or email == "nitin@ee.iitm.ac.in"

# Save session to Supabase
def save_session_to_supabase(email, name, chat_history):
    # Define IST timezone
    ist = pytz.timezone("Asia/Kolkata")
    
    for question, answer in chat_history:
        # Get current datetime in IST and format as "YYYY-MM-DD HH:MM"
        timestamp = datetime.now(ist).strftime("%Y-%m-%d %H:%M")
        
        data = {
            "email": email,
            "name": name if name else None,
            "question": question,
            "answer": answer,
            "timestamp": timestamp,  # Add IST timestamp
        }
        response = supabase.table("chat_sessions_1").insert(data).execute()
        if "error" in response:
            st.error(f"Error saving session data to Supabase: {response['error']['message']}")
            return False
    return True


# Streamlit app
st.title("BDM Chatbot")
st.write("Developed by Lalit & Puneet, students of BS (Applications & Data Science) IIT Madras")
st.write("We encourage you to ask clear and relevant questions to ensure meaningful responses. Remember, Garbage In, Garbage Out.")
st.write("Once you have completed your queries, kindly conclude with the word stop.")
st.write("Disclaimer: All interactions, including questions and answers, are recorded to help us enhance the bot's functionality. By using this chatbot, you agree to the storage of this data.")
st.write("Disclaimer: While this chatbot strives to provide accurate responses, it may occasionally make errors. For verified and comprehensive information, please refer to the official project documents or consult the appropriate team members/mentors.")
# Load model and vector store
model = load_model()
vector_store = reload_vector_store_if_needed()
retrieval_chain = ConversationalRetrievalChain.from_llm(model, retriever=vector_store.as_retriever())

# Initialize session states
if "chat_history" not in st.session_state:
    st.session_state["chat_history"] = []
if "email_validated" not in st.session_state:
    st.session_state["email_validated"] = False
if "session_start_time" not in st.session_state:
    st.session_state["session_start_time"] = datetime.now()

# Email input
email = st.text_input("Enter your email (format: XXfXXXXXXX@ds.study.iitm.ac.in):")
name = st.text_input("Enter your name (optional):")

if email and is_valid_email(email):
    st.session_state["email_validated"] = True
    st.success("Email validated successfully! You can now ask your questions.")
elif email:
    st.error("Invalid email format. Please enter a valid email.")

# Check session time limit
elapsed_time = datetime.now() - st.session_state["session_start_time"]
if elapsed_time > timedelta(minutes=30):
    st.warning("Session time has exceeded 30 minutes. The chatbot will stop now. Refresh to start a new session")
    
    # Save session data to Supabase
    if save_session_to_supabase(email, name, st.session_state["chat_history"]):
        st.success("Session data successfully saved to Supabase!")
    
    # Download session data
    session_data = {
    "email": email,
    "name": name,
    "chat_history": [
        {"question": q, "answer": a, "timestamp": datetime.now().isoformat()}
        for q, a in st.session_state["chat_history"]
    ]
}
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"session_data_{timestamp}.json"
    st.download_button(
        label="Download Session Data",
        data=json.dumps(session_data, indent=4),
        file_name=filename,
        mime="application/json"
    )
    st.stop()


# Chat functionality
if st.session_state["email_validated"]:
    user_input = st.text_input("Pose your Questions:")

    if user_input:
        if user_input.lower() == "stop":
            st.write("Chatbot: Goodbye!")
            session_data = {
                "email": email,
                "name": name,
                "chat_history": st.session_state["chat_history"]
            }

            # Save session data to Supabase
            if save_session_to_supabase(email, name, st.session_state["chat_history"]):
                st.success("Session data successfully saved to Supabase! Refresh to start a new session")

            # Allow download of session data
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"session_data_{timestamp}.json"
            st.download_button(
                label="Download Session Data",
                data=json.dumps(session_data, indent=4),
                file_name=filename,
                mime="application/json"
            )
            st.stop()
        else:
            # Process the question and answer
            response = retrieval_chain.invoke({"question": user_input, "chat_history": st.session_state["chat_history"]})
            answer = response["answer"]
            st.session_state["chat_history"].append((user_input, answer))
            
            # Display the chat history
            st.subheader("Chat History")
            
            # Use a container to make the chat scrollable
            with st.container():
                for i, (question, reply) in enumerate(st.session_state["chat_history"], 1):
                    st.markdown(f"**Q{i}:** {question}")
                    st.markdown(f"**Chatbot:** {reply}")
            
            


